from flask import Flask, render_template, request, send_file, Response
import openai, json
from pptx import Presentation
from pptx.util import Pt
import io

app = Flask(__name__)

openai.api_key = "YOUR_API_KEY"

@app.route('/')
def index():
    return render_template('form.html')

@app.route('/create_ppt', methods=['POST', 'GET'])
def create_ppt():
    try:
        presentation_title = request.form['title']
        print(f"Input Title: {presentation_title}")  # Debugging print

        query_json = """
        {
            "input_text_1": "[[QUERY]]",
            "output_format": "json",
            "json_structure": {
                "slides": "{{presentation_slides}}"
            }
        }
        """

        question = "Generate a 10 slide presentation for the topic. Produce 50 to 60 words per slide. " + presentation_title + ". Each slide should have a {{header}}, {{content}}. The final slide should be a list of discussion questions. Return as JSON."

        prompt = query_json.replace("[[QUERY]]", question)

        completion = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt}])
        response = completion.choices[0].message.content
        r = json.loads(response)

        slide_data = r["slides"]  # Use .get() to safely access dictionary key

        if not slide_data:
            raise ValueError("Invalid response format from OpenAI")

        prs = Presentation()

        # Add title slide with user-inputted title

        for slide in slide_data:
            slide_layout = prs.slide_layouts[1]
            new_slide = prs.slides.add_slide(slide_layout)

            header = slide.get('header')
            content = slide.get('content')

            if header:
                title = new_slide.shapes.title
                title.text = header

            if content:
                shapes = new_slide.shapes
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame

                if isinstance(content, list):
                    content = "\n".join(content)

                tf.text = content

                # Set font manually
                for paragraph in tf.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Calibri'
                        run.font.size = Pt(18)
                        run.font.bold = True

        output_path = "output.pptx"
        prs.save(output_path)


        # Return success message or link to download
        return render_template('form.html', presentation_title=presentation_title)

    except Exception as e:
        print(f"Error creating presentation: {str(e)}")
        return str(e), 500

if __name__ == '__main__':
    app.run(debug=True)
