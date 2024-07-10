from flask import Flask, request, jsonify
import openai
import json
from pptx import Presentation
from pptx.util import Pt

app = Flask(__name__)

# Ganti dengan API key Anda
openai.api_key = "YOUR_API_KEY"

@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    data = request.get_json()
    presentation_title = data.get('title', '')

    query_json = """{
        "input_text": "[[QUERY]]",
        "output_format": "json",
        "json_structure": {
            "slides":"{{presentation_slides}}"
        }
    }"""

    question = f"Generate a 10 slide presentation for the topic. Produce 50 to 60 words per slide. {presentation_title}. Each slide should have a {{header}}, {{content}}. The final slide should be a list of discussion questions. Return as JSON."

    prompt = query_json.replace("[[QUERY]]", question)
    completion = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=[{"role": "user", "content": prompt}])
    response = completion.choices[0].message.content

    r = json.loads(response)
    slide_data = r["slides"]

    prs = Presentation()

    for slide in slide_data:
        slide_layout = prs.slide_layouts[1]
        new_slide = prs.slides.add_slide(slide_layout)

        if slide['header']:
            title = new_slide.shapes.title
            title.text = slide['header']

        if slide['content']:
            shapes = new_slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame

            # Convert list to string if necessary
            if isinstance(slide['content'], list):
                content = "\n".join(slide['content'])
            else:
                content = slide['content']

            tf.text = content

            # Set font manually
            for paragraph in tf.paragraphs:
                for run in paragraph.runs:
                    run.font.name = 'Calibri'
                    run.font.size = Pt(18)
                    run.font.bold = True

    output_path = "output.pptx"
    prs.save(output_path)

    return jsonify({"message": "PPT generated successfully!", "file_path": output_path})

if __name__ == '__main__':
    app.run(debug=True)
